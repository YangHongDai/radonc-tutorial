#!/usr/bin/env python3
"""Import legacy radiation oncology exam questions into radonc-quiz/data.js.

The source quiz folder contains a mixture of Word, PDF, slides, images, notes,
and exam sheets. This importer extracts readable text, detects MCQ-style
questions, keeps source provenance, and emits the simple QUIZ_DATA schema used
by radonc-quiz.
"""

from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from collections import Counter, defaultdict
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
QUIZ_DIR = ROOT / "quiz"
OUT_DATA = ROOT / "radonc-quiz" / "data.js"
REPORT = ROOT / "radonc-quiz" / "IMPORT_REPORT.md"
TMP_TEXT = ROOT / "tmp" / "quiz_text"
PYTHON = Path("/Users/arthurdai/.cache/codex-runtimes/codex-primary-runtime/dependencies/python/bin/python3")
SOFFICE = Path("/Users/arthurdai/.cache/codex-runtimes/codex-primary-runtime/dependencies/bin/soffice")

TEXT_EXTS = {".txt", ".html", ".htm"}
DOC_EXTS = {".doc", ".docx", ".rtf"}
PDF_EXTS = {".pdf"}
PPT_EXTS = {".ppt", ".pptx"}
SHEET_EXTS = {".xlsx", ".xls"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg"}
SKIP_EXTS = {".dcm", ".ds_store", ".tmp", ".zip", ".rar", ".numbers", ".pages", ".gif"}
SKIP_PATH_PARTS = {
    "筆記整理",
    "2020回憶綜整",
}
LIKELY_QUESTION_HINTS = {
    "考題", "Recall", "recall", "題", "密件", "試題", "口試", "board", "exam",
    "Physics", "Biology", "CNS", "PED", "GYN", "GU", "GI", "HN", "Lung",
    "Breast", "NHL", "sarcoma", "Skin", "Palliative",
}

FULLWIDTH = str.maketrans({
    "Ａ": "A", "Ｂ": "B", "Ｃ": "C", "Ｄ": "D",
    "ａ": "A", "ｂ": "B", "ｃ": "C", "ｄ": "D",
    "Ｔ": "T", "Ｆ": "F", "（": "(", "）": ")",
    "：": ":", "．": ".", "、": ".", "　": " ",
})


def run(cmd: list[str], timeout: int = 60) -> subprocess.CompletedProcess:
    return subprocess.run(cmd, text=True, capture_output=True, timeout=timeout)


def clean_text(text: str) -> str:
    text = text.replace("\x00", "")
    text = text.translate(FULLWIDTH)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[\u200b\u200c\u200d\ufeff]", "", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def safe_name(path: Path) -> str:
    rel = path.relative_to(QUIZ_DIR)
    return re.sub(r"[^A-Za-z0-9_.-]+", "_", str(rel)) + ".txt"


def extract_doc(path: Path) -> str:
    cmd = ["textutil", "-convert", "txt", "-stdout", str(path)]
    cp = run(cmd, timeout=90)
    if cp.returncode == 0 and cp.stdout.strip():
        return cp.stdout
    return ""


def extract_pdf(path: Path) -> str:
    name = path.name
    if path.stat().st_size > 8_000_000 and not any(h in name for h in LIKELY_QUESTION_HINTS):
        return ""
    try:
        import pdfplumber  # type: ignore

        chunks: list[str] = []
        with pdfplumber.open(path) as pdf:
            if len(pdf.pages) > 80 and not any(h in name for h in LIKELY_QUESTION_HINTS):
                return ""
            for i, page in enumerate(pdf.pages[:80]):
                chunks.append(f"\n--- page {i + 1} ---\n")
                chunks.append(page.extract_text(x_tolerance=1, y_tolerance=3) or "")
        return "\n".join(chunks)
    except Exception as exc:
        return f"[PDF extraction failed: {exc}]"


def extract_pptx(path: Path) -> str:
    if path.suffix.lower() == ".ppt":
        return ""
    if path.suffix.lower() == ".pptx":
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(path)
            chunks: list[str] = []
            for i, slide in enumerate(prs.slides):
                chunks.append(f"\n--- slide {i + 1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
            return "\n".join(chunks)
        except Exception:
            pass

    if SOFFICE.exists():
        with tempfile.TemporaryDirectory() as td:
            cp = run([
                str(SOFFICE), "--headless", "--convert-to", "txt:Text",
                "--outdir", td, str(path)
            ], timeout=120)
            for txt in Path(td).glob("*.txt"):
                data = txt.read_text(errors="ignore")
                if data.strip():
                    return data
            return cp.stdout + "\n" + cp.stderr
    return ""


def extract_sheet(path: Path) -> str:
    if path.suffix.lower() == ".xlsx":
        try:
            import openpyxl  # type: ignore

            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            chunks: list[str] = []
            for ws in wb.worksheets:
                chunks.append(f"\n--- sheet {ws.title} ---")
                for row in ws.iter_rows(values_only=True):
                    vals = [str(v) for v in row if v is not None and str(v).strip()]
                    if vals:
                        chunks.append("\t".join(vals))
            return "\n".join(chunks)
        except Exception as exc:
            return f"[XLSX extraction failed: {exc}]"
    if SOFFICE.exists():
        with tempfile.TemporaryDirectory() as td:
            cp = run([
                str(SOFFICE), "--headless", "--convert-to", "csv",
                "--outdir", td, str(path)
            ], timeout=120)
            chunks = []
            for csv in Path(td).glob("*.csv"):
                chunks.append(csv.read_text(errors="ignore"))
            return "\n".join(chunks) or cp.stdout + "\n" + cp.stderr
    return ""


def swift_ocr_available() -> bool:
    return bool(shutil.which("swift"))


def extract_image_ocr(path: Path) -> str:
    if os.environ.get("QUIZ_OCR", "0") != "1":
        return ""
    # Restrict OCR to likely question screenshots to keep runtime sane.
    rel = str(path.relative_to(QUIZ_DIR))
    if not any(key in rel for key in ["換題", "PED", "考題", "實際考題", "Recall", "recall"]):
        return ""
    if not swift_ocr_available():
        return ""
    swift = r'''
import Foundation
import Vision
import AppKit

let path = CommandLine.arguments[1]
guard let image = NSImage(contentsOfFile: path),
      let tiff = image.tiffRepresentation,
      let bitmap = NSBitmapImageRep(data: tiff),
      let cgImage = bitmap.cgImage else {
  exit(2)
}
let request = VNRecognizeTextRequest()
request.recognitionLevel = .accurate
request.usesLanguageCorrection = true
request.recognitionLanguages = ["zh-Hant", "zh-Hans", "en-US"]
let handler = VNImageRequestHandler(cgImage: cgImage, options: [:])
try handler.perform([request])
let lines = (request.results ?? []).compactMap { $0.topCandidates(1).first?.string }
print(lines.joined(separator: "\n"))
'''
    with tempfile.NamedTemporaryFile("w", suffix=".swift", delete=False) as fh:
        fh.write(swift)
        script = fh.name
    try:
        cp = run(["swift", script, str(path)], timeout=60)
        return cp.stdout if cp.returncode == 0 else ""
    finally:
        try:
            os.unlink(script)
        except OSError:
            pass


def extract_file(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    try:
        if ext in TEXT_EXTS:
            return path.read_text(errors="ignore"), "text"
        if ext in DOC_EXTS:
            return extract_doc(path), "doc"
        if ext in PDF_EXTS:
            return extract_pdf(path), "pdf"
        if ext in PPT_EXTS:
            return extract_pptx(path), "ppt"
        if ext in SHEET_EXTS:
            return extract_sheet(path), "sheet"
        if ext in IMAGE_EXTS:
            return extract_image_ocr(path), "image-ocr"
    except Exception as exc:
        return f"[Extraction failed: {exc}]", "failed"
    return "", "skipped"


def lines_of(text: str) -> list[str]:
    return [ln.strip() for ln in clean_text(text).splitlines() if ln.strip()]


def is_option_line(line: str) -> re.Match[str] | None:
    return re.match(r"^\s*([ABCD])[\)\.、:：]\s*(.+)$", line.translate(FULLWIDTH), re.I)


def is_answer_line(line: str) -> re.Match[str] | None:
    s = line.translate(FULLWIDTH).strip()
    patterns = [
        r"^\(?\s*\d{1,3}\s*答案\s*\)?\s*[:：]?\s*([ABCDTF])\b(.*)$",
        r"^答案\s*[:：]?\s*([ABCDTF])\b(.*)$",
        r"^(?:Ans|Answer)\s*[:：]?\s*([ABCDTF])\b(.*)$",
        r"^([ABCD])\s*(?:\([^)]{1,30}\))?\s*$",
        r"^([ABCD])\s+[(（][^)）]{1,30}[)）].*$",
    ]
    for pat in patterns:
        m = re.match(pat, s, re.I)
        if m:
            return m
    return None


def is_question_start(line: str) -> re.Match[str] | None:
    s = line.translate(FULLWIDTH)
    return re.match(r"^\s*(\d{1,3})[\.\)]\s*(?:\([^)]*\)\s*)?(.{4,})$", s)


def compact_join(parts: list[str]) -> str:
    return re.sub(r"\s+", " ", " ".join(p.strip() for p in parts if p.strip())).strip()


def tidy_question_text(text: str) -> str:
    text = clean_text(str(text or ""))
    text = re.sub(r"^\s*\d{1,3}\s*[\.\)]\s*(?:\(\s*\)\s*)?", "", text)
    text = re.sub(r"^\s*\(\s*\)\s*", "", text)
    text = re.sub(r"^[-–—]+\s*", "", text)
    text = re.sub(r"\s+([?？,，.。:：;；])", r"\1", text)
    text = re.sub(r"([（(])\s+|\s+([）)])", lambda m: m.group(1) or m.group(2), text)
    return text.strip()


def is_dependent_stem(stem: str) -> bool:
    return bool(re.match(r"^(承接上題|承上題|根據上題|依上題|上題)", tidy_question_text(stem)))


def source_context_before(text: str, stem: str, max_lines: int = 12) -> str:
    ls = lines_of(text)
    target = tidy_question_text(stem)
    target_key = re.sub(r"\s+", "", target[:28])
    if not target_key:
        return ""
    for i, ln in enumerate(ls):
        line_key = re.sub(r"\s+", "", tidy_question_text(ln)[:28])
        if line_key and (line_key.startswith(target_key[:18]) or target_key.startswith(line_key[:18])):
            start = max(0, i - max_lines)
            context_lines = [
                tidy_question_text(x)
                for x in ls[start:i]
                if tidy_question_text(x) and not re.match(r"^(命題者|HYPERLINK)", tidy_question_text(x), re.I)
            ]
            return compact_join(context_lines)
    return ""


def split_exp(block_lines: list[str]) -> str:
    text = "\n".join(block_lines)
    m = re.search(r"(?:註解|解析|說明)\)?\s*[:：]\s*(.+)$", text, re.S)
    if m:
        return compact_join(m.group(1).splitlines())
    # Fall back to the prose after reference/difficulty lines.
    useful = []
    keep = False
    for ln in block_lines:
        if re.search(r"註解|解析|說明|參考文獻|Perez|NCCN|ASTRO|p\d+", ln, re.I):
            keep = True
        if keep and not is_answer_line(ln):
            useful.append(ln)
    return compact_join(useful[-8:])


SUBJECT_RULES = [
    (("淋巴瘤、肉瘤與皮膚", "皮膚癌"), ["皮膚癌", "skin cancer", "basal cell", "squamous cell carcinoma of skin", "melanoma", "merkel", "蕈狀肉芽腫"]),
    (("乳癌與消化道", "乳癌"), ["乳癌", "乳房", "breast", "dcis", "mastectomy", "lumpectomy", "tamoxifen", "her2"]),
    (("乳癌與消化道", "消化道癌症"), ["直腸", "大腸", "結腸", "食道", "胃癌", "胃部", "肝癌", "肝細胞", "膽管", "胰臟", "肛門癌", "rectal", "rectum", "colon", "esophageal", "esophagus", "gastric", "stomach", "hcc", "hepatocellular", "cholangi", "pancrea", "anal cancer"]),
    (("泌尿與婦科", "婦科腫瘤"), ["子宮頸", "子宮內膜", "卵巢", "陰道癌", "外陰", "婦科", "cervix", "cervical", "endometrial", "ovarian", "vulvar", "vagina", "figo"]),
    (("泌尿與婦科", "泌尿生殖癌"), ["攝護腺", "前列腺", "膀胱", "睪丸", "輸尿管", "腎盂", "泌尿", "prostate", "psa", "gleason", "bladder", "testis", "testicular", "ureter", "renal pelvis", "urothelial"]),
    (("頭頸與胸腔", "頭頸部癌症"), ["鼻咽", "口咽", "下咽", "喉癌", "聲門", "唾液腺", "甲狀腺", "頭頸", "舌癌", "口腔癌", "嗅神經母細胞瘤", "蝶竇", "鼻腔", "副鼻竇", "nasophary", "npc", "orophary", "hypophary", "laryn", "salivary", "thyroid", "head and neck", "oral cavity", "esthesioneuroblastoma", "sphenoid sinus", "nasal cavity", "paranasal"]),
    (("頭頸與胸腔", "肺癌"), ["肺癌", "小細胞肺癌", "非小細胞", "縱膈", "胸腔", "胸膜", "胸腺", "間皮瘤", "nsclc", "sclc", "lung cancer", "mediastinum", "mediastinal", "thoracic", "thymoma", "thymic", "mesothelioma", "pneumonitis", "great vessels", "extragonadal germ cell"]),
    (("中樞神經系統", "原發性 CNS 腫瘤"), ["中樞神經", "腦瘤", "腦膜瘤", "膠質", "神經膠", "髓母", "室管膜", "glioma", "glioblastoma", "gbm", "meningioma", "medulloblastoma", "ependymoma"]),
    (("兒童與緩和", "兒童放射腫瘤"), ["小兒", "兒童", "pediatric", "paediatric", "wilms", "neuroblastoma", "rhabdomyosarcoma", "ewing"]),
    (("兒童與緩和", "緩和放射治療"), ["緩和", "骨轉移", "腦轉移", "脊髓壓迫", "上腔靜脈", "止血", "palliative", "bone metast", "brain metast", "cord compression", "svc syndrome", "hemostatic"]),
    (("淋巴瘤、肉瘤與皮膚", "淋巴瘤"), ["淋巴瘤", "何杰金", "霍奇金", "hodgkin", "lymphoma", "nhl", "dlbcl", "malt"]),
    (("淋巴瘤、肉瘤與皮膚", "軟組織肉瘤"), ["肉瘤", "soft tissue sarcoma", "sarcoma", "sts", "liposarcoma", "leiomyosarcoma"]),
    (("基礎入門", "放射生物學"), ["放射生物", "radiobiology", "radiobiologic", "dna", "染色體", "細胞週期", "細胞傷害", "修復", "凋亡", "存活曲線", "決定性效應", "機率性效應", "白內障", "生物效應", "oer", "let", "rbe", "4r", "hypoxia", "sublethal", "repair", "oxygen enhancement", "radiation weighting factor", "weighting factor", "linear quadratic", "alpha beta", "α/β"]),
    (("基礎入門", "放射物理學"), ["放射物理", "physics", "pdd", "tmr", "tar", "ssd", "sad", "電子束", "光子", "質子", "bragg", "bolus", "楔形板", "半影", "isocenter", "monitor unit", "ctdi", "輻射防護", "半衰期", "衰變常數"]),
]


PATH_FALLBACK_RULES = [
    (("基礎入門", "放射物理學"), ["物理", "physics", "輻防", "放射師"]),
    (("基礎入門", "放射生物學"), ["生物", "biology", "radiobiology"]),
    (("中樞神經系統", "原發性 CNS 腫瘤"), ["cns"]),
    (("頭頸與胸腔", "頭頸部癌症"), ["hn", "head", "neck", "頭頸"]),
    (("頭頸與胸腔", "肺癌"), ["lung", "chest"]),
    (("乳癌與消化道", "乳癌"), ["breast"]),
    (("乳癌與消化道", "消化道癌症"), ["gi"]),
    (("泌尿與婦科", "婦科腫瘤"), ["gyn"]),
    (("泌尿與婦科", "泌尿生殖癌"), ["gu", "uro"]),
    (("兒童與緩和", "兒童放射腫瘤"), ["ped", "pediatric"]),
    (("淋巴瘤、肉瘤與皮膚", "淋巴瘤"), ["hema", "lymph"]),
]


def score_keywords(text: str, keywords: list[str]) -> int:
    score = 0
    low = text.lower()
    for kw in keywords:
        k = kw.lower()
        if re.fullmatch(r"[a-z0-9]{1,3}", k):
            matched = bool(re.search(rf"(?<![a-z0-9]){re.escape(k)}(?![a-z0-9])", low))
        else:
            matched = k in low
        if matched:
            score += 3 if len(k) >= 4 else 2
    return score


def infer_subject(q: dict) -> tuple[str, str]:
    stem = tidy_question_text(str(q.get("stem", ""))).lower()
    rest = " ".join(str(q.get(k, "")) for k in ["A", "B", "C", "D", "exp"])
    rest = tidy_question_text(rest).lower()
    scores = []
    for dest, keywords in SUBJECT_RULES:
        s = score_keywords(stem, keywords) * 4 + score_keywords(rest, keywords)
        if s:
            scores.append((s, dest))
    if scores:
        scores.sort(key=lambda item: item[0], reverse=True)
        return scores[0][1]

    rel = str(q.get("source", "")).lower()
    for dest, keywords in PATH_FALLBACK_RULES:
        if any(k.lower() in rel for k in keywords):
            return dest
    if "口試" in rel:
        return ("口試複習", "口試題庫")
    return ("綜合題庫", "未分類考題")


def parse_numbered(path: Path, text: str) -> list[dict]:
    ls = lines_of(text)
    starts = [(i, is_question_start(ln)) for i, ln in enumerate(ls) if is_question_start(ln)]
    out = []
    for idx, (start_i, m) in enumerate(starts):
        end_i = starts[idx + 1][0] if idx + 1 < len(starts) else len(ls)
        block = ls[start_i:end_i]
        if not m:
            continue
        stem_parts = [m.group(2).strip()]
        options: dict[str, list[str]] = {}
        answer = None
        answer_tail = ""
        current = None
        option_started = False
        after_answer: list[str] = []
        for ln in block[1:]:
            am = is_answer_line(ln)
            if am and re.search(r"答案|Ans|Answer", ln, re.I):
                answer = am.group(1).upper()
                answer_tail = am.group(2) if am.lastindex and am.lastindex >= 2 else ""
                current = None
                continue
            om = is_option_line(ln)
            if om and answer is None:
                current = om.group(1).upper()
                options[current] = [om.group(2).strip()]
                option_started = True
                continue
            if answer is not None:
                after_answer.append(ln)
            elif current and option_started:
                if re.search(r"^\(?\s*\d{1,3}\s*(答案|難易度|參考文獻|註解)", ln):
                    current = None
                else:
                    options[current].append(ln)
            elif not option_started:
                stem_parts.append(ln)
        if not answer:
            for ln in block:
                am = is_answer_line(ln)
                if am:
                    answer = am.group(1).upper()
                    answer_tail = am.group(2) if am.lastindex and am.lastindex >= 2 else ""
                    break
        if answer and len(options) >= 2 and answer in "ABCD":
            stem = tidy_question_text(compact_join(stem_parts))
            q = {
                "type": "single",
                "stem": stem,
                "ans": answer,
                "exp": split_exp(after_answer) or f"正解為 {answer}。原始題庫未提供可辨識的詳解，請搭配來源檔複核。",
                "source": str(path.relative_to(QUIZ_DIR)),
            }
            for k in "ABCD":
                if k in options:
                    q[k] = tidy_question_text(compact_join(options[k]))
            if q.get(answer):
                out.append(q)
    return out


def parse_loose(path: Path, text: str) -> list[dict]:
    ls = lines_of(text)
    out = []
    i = 0
    question_keywords = re.compile(r"(何者|下列|有關|關於|請問|哪一|何項)")
    while i < len(ls) - 5:
        line = ls[i]
        if question_keywords.search(line) and not is_option_line(line) and not is_question_start(line):
            # Try stem + four option lines + answer line.
            opts = ls[i + 1:i + 5]
            ans_line = ls[i + 5] if i + 5 < len(ls) else ""
            am = is_answer_line(ans_line)
            if am and all(not is_answer_line(x) for x in opts):
                answer = am.group(1).upper()
                answer_tail = am.group(2) if am.lastindex and am.lastindex >= 2 else ""
                if answer in "ABCD":
                    exp_lines = []
                    j = i + 6
                    while j < len(ls) and len(exp_lines) < 8:
                        if question_keywords.search(ls[j]) and j + 5 < len(ls) and is_answer_line(ls[j + 5]):
                            break
                        exp_lines.append(ls[j])
                        j += 1
                    q = {
                        "type": "single",
                        "stem": tidy_question_text(compact_join([line])),
                        "A": tidy_question_text(compact_join([opts[0]])),
                        "B": tidy_question_text(compact_join([opts[1]])),
                        "C": tidy_question_text(compact_join([opts[2]])),
                        "D": tidy_question_text(compact_join([opts[3]])),
                        "ans": answer,
                        "exp": compact_join(exp_lines) or f"正解為 {answer}。原始題庫未提供可辨識的詳解，請搭配來源檔複核。",
                        "source": str(path.relative_to(QUIZ_DIR)),
                    }
                    out.append(q)
                    i = max(j, i + 6)
                    continue
        i += 1
    return out


def normalize_answer(q: dict) -> dict | None:
    ans = str(q.get("ans", "")).strip().upper().translate(FULLWIDTH)
    if ans in {"TRUE", "T"}:
        return {"type": "tf", "stem": q["stem"], "ans": "T", "exp": q.get("exp", ""), "source": q.get("source", "")}
    if ans in {"FALSE", "F"}:
        return {"type": "tf", "stem": q["stem"], "ans": "F", "exp": q.get("exp", ""), "source": q.get("source", "")}
    if ans not in "ABCD":
        return None
    q["ans"] = ans
    if ans not in q:
        return None
    if q.get("type") == "single" and any(k not in q or not str(q.get(k, "")).strip() for k in "ABCD"):
        return None
    for k in list("ABCD"):
        if k in q:
            q[k] = tidy_question_text(str(q[k]))
    q["stem"] = tidy_question_text(str(q["stem"]))
    q["exp"] = tidy_question_text(str(q.get("exp", ""))) or f"正解為 {ans}。"
    return q


def extract_all() -> tuple[list[tuple[Path, str, str]], Counter]:
    TMP_TEXT.mkdir(parents=True, exist_ok=True)
    extracted: list[tuple[Path, str, str]] = []
    stats: Counter = Counter()
    files = [p for p in sorted(QUIZ_DIR.rglob("*")) if p.is_file()]
    total = len(files)
    for n, path in enumerate(files, 1):
        if n == 1 or n % 10 == 0:
            print(f"[extract] {n}/{total} {path.relative_to(QUIZ_DIR)}", flush=True)
        if path.name.startswith("~$") or path.name.startswith("."):
            stats["skipped_temp"] += 1
            continue
        rel_parts = set(path.relative_to(QUIZ_DIR).parts)
        if rel_parts & SKIP_PATH_PARTS:
            stats["skipped_note_folder"] += 1
            continue
        ext = path.suffix.lower()
        if ext in SKIP_EXTS:
            stats[f"skipped_{ext or 'none'}"] += 1
            continue
        text, kind = extract_file(path)
        text = clean_text(text)
        if text:
            (TMP_TEXT / safe_name(path)).write_text(text, encoding="utf-8")
            extracted.append((path, text, kind))
            stats[f"extracted_{kind}"] += 1
        else:
            stats[f"empty_{kind}"] += 1
    return extracted, stats


def build_data(questions: list[dict]) -> dict:
    subjects: dict[str, dict[str, list[dict]]] = defaultdict(lambda: defaultdict(list))
    for q in questions:
        clean_q = {k: v for k, v in q.items() if k in {"type", "stem", "A", "B", "C", "D", "ans", "exp", "source"}}
        top, sub = infer_subject(clean_q)
        subjects[top][sub].append(clean_q)
    return {
        "radonc": {
            "label": "☢️ OncoBeam 放射腫瘤學考題",
            "subjects": {top: dict(subs) for top, subs in subjects.items()},
        }
    }


def main() -> int:
    extracted, stats = extract_all()
    all_q: list[dict] = []
    per_file: Counter = Counter()
    missing_answer_candidates = 0
    for path, text, _kind in extracted:
        qs = parse_numbered(path, text)
        seen = {(q["stem"], q.get("A", ""), q.get("B", ""), q.get("C", ""), q.get("D", "")) for q in qs}
        loose = []
        for q in parse_loose(path, text):
            key = (q["stem"], q.get("A", ""), q.get("B", ""), q.get("C", ""), q.get("D", ""))
            if key not in seen:
                loose.append(q)
        qs.extend(loose)
        previous_context = ""
        for q in qs:
            if previous_context and is_dependent_stem(q.get("stem", "")):
                q["stem"] = f"題組背景：{previous_context}\n\n{q['stem']}"
            elif is_dependent_stem(q.get("stem", "")):
                fallback_context = source_context_before(text, q.get("stem", ""))
                if fallback_context:
                    q["stem"] = f"題組背景：{fallback_context}\n\n{q['stem']}"
            nq = normalize_answer(q)
            if nq:
                all_q.append(nq)
                per_file[str(path.relative_to(QUIZ_DIR))] += 1
                previous_context = nq["stem"]
            else:
                missing_answer_candidates += 1

    # Deduplicate identical questions/options while preserving first source.
    deduped = []
    seen_keys = set()
    duplicates = 0
    for q in all_q:
        key = (
            re.sub(r"\s+", "", tidy_question_text(q.get("stem", ""))),
            re.sub(r"\s+", "", tidy_question_text(q.get("A", ""))),
            re.sub(r"\s+", "", tidy_question_text(q.get("B", ""))),
            re.sub(r"\s+", "", tidy_question_text(q.get("C", ""))),
            re.sub(r"\s+", "", tidy_question_text(q.get("D", ""))),
            q.get("ans", ""),
        )
        if key in seen_keys:
            duplicates += 1
            continue
        seen_keys.add(key)
        deduped.append(q)

    data = build_data(deduped)
    OUT_DATA.write_text(
        "/* Generated from quiz/ by scripts/import_quiz_bank.py. Chinese-only question bank. */\n"
        "window.QUIZ_DATA = "
        + json.dumps(data, ensure_ascii=False, indent=2)
        + ";\n",
        encoding="utf-8",
    )

    section_counts = []
    for top, subs in data["radonc"]["subjects"].items():
        for sub, arr in subs.items():
            section_counts.append((top, sub, len(arr)))
    section_counts.sort()

    report = [
        "# Quiz Import Report",
        "",
        f"- Extracted readable files: {len(extracted)}",
        f"- Imported questions: {len(deduped)}",
        f"- Duplicate questions removed: {duplicates}",
        f"- Parsed candidates discarded for invalid/missing answers: {missing_answer_candidates}",
        "",
        "## Extraction Stats",
        "",
    ]
    for key, val in sorted(stats.items()):
        report.append(f"- {key}: {val}")
    report += ["", "## Section Counts", ""]
    for top, sub, n in section_counts:
        report.append(f"- {top} / {sub}: {n}")
    report += ["", "## Source Files With Imported Questions", ""]
    for file, n in per_file.most_common():
        report.append(f"- `{file}`: {n}")
    REPORT.write_text("\n".join(report) + "\n", encoding="utf-8")
    print(f"Imported {len(deduped)} questions into {OUT_DATA}")
    print(f"Report written to {REPORT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
